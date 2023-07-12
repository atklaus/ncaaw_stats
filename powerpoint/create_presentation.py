from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title, content, bullet=False):
    slide_layout = prs.slide_layouts[5] if bullet else prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Create a text box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txBox.text_frame

    for line in content.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        if bullet:
            p.level = 0
        else:
            p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # ... (The rest of the script remains the same)

    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Predicting NBA Success and Performance: A Literature Review"
    subtitle.text = "A Review of Three Research Papers"

    # Slide 1: Overview
    content = '''Three research papers examined:
  1. Kannan et al. (2018): A Machine Learning Approach
  2. Maymin (2021): Using Scouting Reports Text
  3. Rodenberg & Kim (2011): Precocity and Labor Market Outcomes'''
    add_slide(prs, "Overview", content, bullet=True)

    # Slide 2-4: Research Papers
    papers = [
        {
            'title': 'Predicting National Basketball Association Success: A Machine Learning Approach - Kannan et al. (2018)',
            'content': '''Data Sources: * Historical NBA player data (1980-2017) from Basketball-Reference.com
      * NBA combine data from NBA.com
      * College statistics from Sports-Reference.com

    NBA Success Definition: * Player Efficiency Rating (PER): A measure of per-minute production standardized such that the league average is 15.
      * Win Shares (WS): An estimate of the number of wins contributed by a player.
      * Win Shares per 48 minutes (WS/48): Win Shares per 48 minutes played.

    Key Findings: * Machine learning model effectively predicted NBA success using college stats and NBA combine data
      * Importance of athleticism and college performance

    Link: https://scholar.smu.edu/datasciencereview/vol1/iss3/7/'''
        },
        {
            'title': 'Using Scouting Reports Text To Predict NCAA → NBA Performance - Maymin (2021)',
            'content': '''Data Sources: * Scouting reports from DraftExpress.com and ESPN.com
      * NCAA and NBA statistics from Sports-Reference.com

    NBA Performance Definition: * Box Plus/Minus (BPM): A box score estimate of the points per 100 possessions a player contributed above a league-average player, translated to an average team.
      * Value Over Replacement Player (VORP): A box score estimate of the points per 100 team possessions a player contributed above a replacement-level (-2.0) player, translated to an average team, and prorated to an 82-game season.

    Key Findings: * Text-based scouting reports contain valuable predictive information
      * Analysis of report content can help predict player performance in the NBA

    Link: https://doi.org/10.1080/2573234X.2021.1873077'''
        },
        {
            'title': 'Precocity and Labor Market Outcomes: Evidence from Professional Basketball - Rodenberg & Kim (2011)',
            'content': '''Data Sources: * NBA player data (1977-2010) from Basketball-Reference.com
      * Player biographical data from various sources

    NBA Success Definition: * Annual earnings: The annual salary of a player.
      * Career length: The number of years a player remains active in the NBA.

    Key Findings: * Precocious athletes (young players who enter the NBA early) earn more and have longer careers
      * Age at entry positively correlated with career earnings and longevity

    Link: https://ssrn.com/abstract=1869144'''
        }
    ]

    for paper in papers:
        add_slide(prs, paper['title'], paper['content'], bullet=True)

    content = '''Diverse definitions of NBA success/performance
Multiple factors can predict player success, including college stats, combine data, scouting reports, and age at entry
Machine learning and textual analysis can help uncover predictive relationships'''
    add_slide(prs, "Conclusion", content, bullet=True)

    # Save the presentation
    prs.save('Predicting_NBA_Success_and_Performance.pptx')

if __name__ == "__main__":
  main()